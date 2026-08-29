#!/usr/bin/env python3
"""Build pitch/GrantOnce.pptx from the sparse 12-slide outline."""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "GrantOnce.pptx"

PAPER = RGBColor(0xF4, 0xEF, 0xE4)
INK = RGBColor(0x1C, 0x16, 0x12)
SEAL = RGBColor(0x9F, 0x1D, 0x1D)
BOX = RGBColor(0x1F, 0x4A, 0x3D)
MUTED = RGBColor(0x6B, 0x62, 0x58)
RULE = RGBColor(0xD9, 0xD0, 0xC0)
WHITE = RGBColor(0xFB, 0xF7, 0xEE)

W = Inches(13.333)
H = Inches(7.5)
FONT = "Microsoft JhengHei"


def set_run_font(run, name: str = FONT, east: str = FONT) -> None:
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    rFonts = rPr.find(qn("a:rFonts"))
    if rFonts is None:
        rFonts = etree.SubElement(rPr, qn("a:rFonts"))
    rFonts.set("ascii", name)
    rFonts.set("hAnsi", name)
    rFonts.set("eastAsia", east)
    rFonts.set("cs", east)


def fill_slide(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_text(
    slide,
    left,
    top,
    width,
    height,
    lines: list[tuple[str, int, RGBColor, bool]],
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(anchor, "t"))
    except Exception:
        pass
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        set_run_font(run)


def footer(slide, n: int) -> None:
    add_text(
        slide,
        Inches(0.6),
        Inches(7.05),
        Inches(8),
        Inches(0.3),
        [("GrantOnce 分匣授權", 11, MUTED, False)],
    )
    add_text(
        slide,
        Inches(11.2),
        Inches(7.05),
        Inches(1.5),
        Inches(0.3),
        [(f"{n:02d} / 12", 11, MUTED, False)],
        align=PP_ALIGN.RIGHT,
    )


def rule(slide, top) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(12.1), Emu(12700))
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = RULE


def kicker(slide, text: str, color: RGBColor = SEAL) -> None:
    add_text(slide, Inches(0.6), Inches(0.38), Inches(12), Inches(0.35), [(text, 12, color, True)])


def build() -> None:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # 01 title
    s = prs.slides.add_slide(blank)
    fill_slide(s, PAPER)
    add_text(s, Inches(0.8), Inches(1.9), Inches(11.5), Inches(1.1), [("分匣授權", 54, INK, True)])
    add_text(s, Inches(0.8), Inches(3.05), Inches(11.5), Inches(0.6), [("GrantOnce", 28, MUTED, False)])
    rule(s, Inches(3.85))
    add_text(
        s,
        Inches(0.8),
        Inches(4.1),
        Inches(11.5),
        Inches(0.7),
        [("只准這一次，而且只准這一匣。", 28, SEAL, True)],
    )
    add_text(
        s,
        Inches(0.8),
        Inches(6.4),
        Inches(11.5),
        Inches(0.4),
        [("週末演示 · 合成資料 · 非正式 MyData", 14, MUTED, False)],
    )
    footer(s, 1)

    # 02 problem
    s = prs.slides.add_slide(blank)
    fill_slide(s, PAPER)
    kicker(s, "問題")
    add_text(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1), [("跨部門補助。", 40, INK, True)])
    add_text(s, Inches(0.8), Inches(2.7), Inches(11.5), Inches(1), [("用戶自己當整合層。", 40, INK, True)])
    add_text(s, Inches(0.8), Inches(4.3), Inches(11.5), Inches(1), [("跑窗口。交影本。對欄位。", 22, MUTED, False)])
    footer(s, 2)

    # 03 wrong solution
    s = prs.slides.add_slide(blank)
    fill_slide(s, PAPER)
    kicker(s, "錯的解法")
    add_text(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1), [("給代理人一張胖 token。", 36, INK, True)])
    add_text(s, Inches(0.8), Inches(2.9), Inches(11.5), Inches(0.7), [("fields:*", 32, SEAL, True)])
    add_text(s, Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.8), [("每個機關都看到全部。", 24, MUTED, False)])
    footer(s, 3)

    # 04 story
    s = prs.slides.add_slide(blank)
    fill_slide(s, PAPER)
    kicker(s, "故事")
    add_text(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.6), [("林曉晴", 36, INK, True)])
    add_text(s, Inches(0.8), Inches(2.1), Inches(11.5), Inches(0.45), [("北市 → 新北　　一歲幼兒　　合成資料", 18, MUTED, False)])
    add_text(
        s,
        Inches(0.8),
        Inches(3.0),
        Inches(11.5),
        Inches(0.8),
        [("「我剛搬家，看我能申請什麼。」", 28, INK, True)],
    )
    add_text(s, Inches(0.8), Inches(4.4), Inches(5.5), Inches(1.2), [("育兒津貼", 28, BOX, True), ("社會局　匣 G-甲", 16, MUTED, False)])
    add_text(s, Inches(7.0), Inches(4.4), Inches(5.5), Inches(1.2), [("冷氣汰換", 28, BOX, True), ("經濟部 × 台電　匣 G-乙", 16, MUTED, False)])
    footer(s, 4)

    # 05 two envelopes
    s = prs.slides.add_slide(blank)
    fill_slide(s, PAPER)
    kicker(s, "兩匣")
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.4), Inches(5.7), Inches(3.6))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RULE
    add_text(s, Inches(0.95), Inches(1.6), Inches(5.2), Inches(3.1), [
        ("G-甲", 16, MUTED, True),
        ("戶籍 + 親子", 28, INK, True),
        ("社會局看得到這些。", 16, MUTED, False),
        ("看不到電號、所得。", 16, MUTED, False),
    ])
    card2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.4), Inches(5.7), Inches(3.6))
    card2.fill.solid()
    card2.fill.fore_color.rgb = WHITE
    card2.line.color.rgb = RULE
    add_text(s, Inches(7.15), Inches(1.6), Inches(5.2), Inches(3.1), [
        ("G-乙", 16, MUTED, True),
        ("電號 + 三月用電", 28, INK, True),
        ("經濟部 × 台電只拿到這些。", 16, MUTED, False),
        ("看不到戶籍、所得。", 16, MUTED, False),
    ])
    add_text(s, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.8), [("所得在金庫。不給。", 28, SEAL, True)])
    footer(s, 5)

    # 06 fail-closed
    s = prs.slides.add_slide(blank)
    fill_slide(s, PAPER)
    kicker(s, "越權關閉")
    add_text(s, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.7), [("乙要戶籍。用匣號當票。", 32, INK, True)])
    add_text(s, Inches(0.8), Inches(2.7), Inches(11.5), Inches(1.2), [("403", 72, SEAL, True)])
    add_text(s, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.7), [("OVERSCOPED · BAD_TICKET。不回傳半包。", 20, MUTED, False)])
    footer(s, 6)

    # 07 revoke
    s = prs.slides.add_slide(blank)
    fill_slide(s, PAPER)
    kicker(s, "送件即撤")
    add_text(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.8), [("送件  →  匣耗用。明文改收據。", 32, INK, True)])
    add_text(s, Inches(0.8), Inches(3.1), Inches(11.5), Inches(0.8), [("重放擷取  →  403。", 36, INK, True)])
    add_text(s, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.6), [("只准這一次。", 22, SEAL, False)])
    footer(s, 7)

    # 08 architecture
    s = prs.slides.add_slide(blank)
    fill_slide(s, PAPER)
    kicker(s, "架構")
    rows = [
        ("HMAC ticket", "runtime 驗 iss／aud／fields／exp。匣號不是憑證。"),
        ("孤立匣", "甲、乙各存各的。拿錯匣 403。"),
        ("規則引擎", "搬家 + 0–2 歲 → 津貼。有電表 → 汰換。"),
        ("LLM 不管授權", "模型不能發匣、不能核准、不能放行欄位。"),
    ]
    y = 1.35
    for title, body in rows:
        add_text(s, Inches(0.8), Inches(y), Inches(12), Inches(0.4), [(title, 22, INK, True)])
        add_text(s, Inches(0.8), Inches(y + 0.38), Inches(12), Inches(0.4), [(body, 16, MUTED, False)])
        y += 1.15
    footer(s, 8)

    # 09 live demo — callouts only; no screenshot (start-screen shot was misleading)
    s = prs.slides.add_slide(blank)
    fill_slide(s, PAPER)
    kicker(s, "實機")
    add_text(
        s,
        Inches(0.6),
        Inches(0.85),
        Inches(12.1),
        Inches(0.4),
        [("三欄同一畫面。看這四件事。", 16, MUTED, False)],
    )
    cols = [
        ("甲匣", "戶籍＋親子", "沒有電號、所得"),
        ("乙匣", "電號＋三月用電", "沒有戶籍、所得"),
        ("越權", "乙要戶籍", "403　半包也不給"),
        ("稽核", "核准／擷取／送件", "撤銷／拒絕"),
    ]
    x = 0.6
    for title, line, note in cols:
        card = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(2.9), Inches(4.4)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RULE
        add_text(
            s,
            Inches(x + 0.18),
            Inches(1.7),
            Inches(2.55),
            Inches(4.0),
            [
                (title, 14, MUTED, True),
                (line, 22, INK, True),
                (note, 14, SEAL if "403" in note else MUTED, False),
            ],
        )
        x += 3.15
    footer(s, 9)

    # 10 MPA
    s = prs.slides.add_slide(blank)
    fill_slide(s, PAPER)
    kicker(s, "MPA")
    add_text(
        s,
        Inches(0.8),
        Inches(1.15),
        Inches(11.5),
        Inches(0.7),
        [("多個委託人，約束同一個代理人。一人一匣。", 24, INK, True)],
    )
    add_text(
        s,
        Inches(0.8),
        Inches(2.15),
        Inches(11.5),
        Inches(0.5),
        [("同一套 Grant。換簽發人。", 28, BOX, True)],
    )
    add_text(
        s,
        Inches(0.8),
        Inches(3.0),
        Inches(11.5),
        Inches(0.55),
        [("主家搬家：簽發人 = 本人（同意）", 20, INK, False)],
    )
    add_text(
        s,
        Inches(0.8),
        Inches(3.55),
        Inches(11.5),
        Inches(0.55),
        [("執法調閱：簽發人 = 法院／搜索票，不是機關甲", 20, INK, False)],
    )
    add_text(
        s,
        Inches(0.8),
        Inches(4.5),
        Inches(11.5),
        Inches(0.9),
        [("甲是收件人。乙仍簽署這包是他的。人還在場。", 18, MUTED, False)],
    )
    footer(s, 10)

    # 11 RBA bonus
    s = prs.slides.add_slide(blank)
    fill_slide(s, PAPER)
    kicker(s, "附頁 · 不是本週演示")
    add_text(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.8), [("RBA", 20, MUTED, True)])
    add_text(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1), [("每個環節出憑證。", 36, INK, True)])
    add_text(s, Inches(0.8), Inches(3.8), Inches(11.5), Inches(0.8), [("供應鏈上的同一句話。", 22, MUTED, False)])
    add_text(s, Inches(0.8), Inches(5.1), Inches(11.5), Inches(0.5), [("本週末不實作。", 16, SEAL, False)])
    footer(s, 11)

    # 12 next
    s = prs.slides.add_slide(blank)
    fill_slide(s, PAPER)
    kicker(s, "下一步")
    add_text(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.8), [("真 MyData 要數發部函。", 32, INK, True)])
    add_text(s, Inches(0.8), Inches(2.8), Inches(11.5), Inches(0.7), [("OID4VP／登入是下一層。", 32, INK, True)])
    add_text(s, Inches(0.8), Inches(3.6), Inches(11.5), Inches(0.5), [("週末用合成資料。", 20, MUTED, False)])
    rule(s, Inches(4.0))
    add_text(
        s,
        Inches(0.8),
        Inches(4.25),
        Inches(11.5),
        Inches(0.8),
        [("只准這一次，而且只准這一匣。", 28, SEAL, True)],
    )
    footer(s, 12)

    prs.save(OUT)
    print(f"wrote {OUT}")


if __name__ == "__main__":
    build()
